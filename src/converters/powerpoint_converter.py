"""
Конвертер PowerPoint презентаций в Markdown
"""
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from loguru import logger
from .base import FileConverter


class PowerPointConverter(FileConverter):
    """
    Конвертер для PowerPoint презентаций
    Поддерживает: .pptx, .ppt (через LibreOffice)
    """

    def __init__(self):
        super().__init__(['.pptx', '.ppt'])
        self._soffice_cmd = None  # Будет установлено в _check_libreoffice
        self.has_pptx = self._check_pptx()
        self.has_libreoffice = self._check_libreoffice()

        if not self.has_pptx and not self.has_libreoffice:
            logger.warning("Ни python-pptx, ни LibreOffice не найдены. Конвертация PowerPoint файлов недоступна.")

    def _check_pptx(self) -> bool:
        """Проверяет доступность python-pptx"""
        try:
            import pptx
            return True
        except ImportError:
            return False

    def _check_libreoffice(self) -> bool:
        """Проверяет наличие LibreOffice в системе"""
        import platform

        # Проверка для Windows
        if platform.system() == 'Windows':
            # Стандартные пути установки LibreOffice в Windows
            possible_paths = [
                r'C:\Program Files\LibreOffice\program\soffice.exe',
                r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            ]
            for path in possible_paths:
                if os.path.exists(path):
                    self._soffice_cmd = path
                    logger.debug(f"LibreOffice найден: {path}")
                    return True

        # Проверка для Linux/macOS через PATH
        commands = ['soffice', 'libreoffice', 'loffice']
        for cmd in commands:
            try:
                result = subprocess.run(
                    [cmd, '--version'],
                    capture_output=True,
                    timeout=5
                )
                if result.returncode == 0:
                    self._soffice_cmd = cmd
                    logger.debug(f"LibreOffice найден: {cmd}")
                    return True
            except (subprocess.TimeoutExpired, FileNotFoundError):
                continue

        return False

    def convert(self, input_path: Path, output_path: Path) -> bool:
        """
        Конвертирует PowerPoint в Markdown

        :param input_path: Путь к исходному файлу
        :param output_path: Путь для сохранения markdown файла
        :return: True если конвертация успешна
        """
        if not self.has_pptx:
            logger.warning(f"python-pptx не установлен. Установите: pip install python-pptx")
            return False

        # Для .ppt используем LibreOffice для конвертации в .pptx
        if input_path.suffix.lower() == '.ppt':
            if self.has_libreoffice:
                return self._convert_ppt_with_libreoffice(input_path, output_path)
            else:
                logger.warning(f"LibreOffice не найден. Невозможно конвертировать старый формат .ppt: {input_path.name}")
                logger.info(f"Установите LibreOffice для поддержки .ppt файлов: https://www.libreoffice.org/download/")
                return False

        # Для .pptx используем python-pptx
        try:
            return self._convert_with_pptx(input_path, output_path)
        except Exception as e:
            logger.error(f"Ошибка при конвертации {input_path}: {e}")
            return False

    def _convert_with_pptx(self, input_path: Path, output_path: Path) -> bool:
        """Конвертирует с помощью python-pptx"""
        from pptx import Presentation

        try:
            # Открываем презентацию
            prs = Presentation(input_path)

            # Создаем директорию для выходного файла
            output_path.parent.mkdir(parents=True, exist_ok=True)

            with open(output_path, 'w', encoding='utf-8') as f:
                # YAML заголовок
                f.write("---\n")
                f.write(f"source_file: {input_path.name}\n")
                f.write(f"original_format: {input_path.suffix.lower()}\n")
                f.write(f"total_slides: {len(prs.slides)}\n")
                f.write(f"converted_by: PowerPointConverter\n")
                f.write("---\n\n")

                # Заголовок
                f.write(f"# {input_path.stem}\n\n")
                f.write(f"**Тип:** PowerPoint презентация\n")
                f.write(f"**Слайдов:** {len(prs.slides)}\n\n")

                # Обрабатываем каждый слайд
                for i, slide in enumerate(prs.slides, 1):
                    f.write(f"\n## Слайд {i}\n\n")

                    # Извлекаем текст из слайда
                    slide_text = []

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text = shape.text.strip()
                            if text:
                                slide_text.append(text)

                    if slide_text:
                        for text in slide_text:
                            # Определяем, является ли текст заголовком (обычно первый текст)
                            if text == slide_text[0] and len(text) < 100:
                                f.write(f"### {text}\n\n")
                            else:
                                f.write(f"{text}\n\n")
                    else:
                        f.write("*(Слайд без текста)*\n\n")

                    # Если есть заметки докладчика
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            f.write(f"**Заметки докладчика:**\n\n")
                            f.write(f"> {notes}\n\n")

            return True

        except Exception as e:
            logger.error(f"Ошибка python-pptx: {e}")
            return False

    def _convert_ppt_with_libreoffice(self, input_path: Path, output_path: Path) -> bool:
        """
        Конвертация .ppt файлов с использованием LibreOffice

        Этапы:
        1. Конвертируем .ppt → .pptx через LibreOffice
        2. Конвертируем .pptx → .md через python-pptx

        :param input_path: Путь к .ppt файлу
        :param output_path: Путь для сохранения .md файла
        :return: True если успешно
        """
        # Создаем временную директорию для промежуточных файлов
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                temp_dir_path = Path(temp_dir)

                # Используем сохраненную команду LibreOffice
                soffice_cmd = self._soffice_cmd
                if not soffice_cmd:
                    logger.error("LibreOffice не найден")
                    return False

                # Шаг 1: Конвертируем .ppt в .pptx через LibreOffice
                logger.debug(f"Конвертация .ppt → .pptx через LibreOffice: {input_path.name}")

                # Копируем файл во временную директорию (LibreOffice имеет проблемы с некоторыми путями)
                temp_input = temp_dir_path / input_path.name
                shutil.copy2(input_path, temp_input)

                # Создаем временный user profile для LibreOffice (решает проблемы с блокировками)
                user_profile_dir = temp_dir_path / 'libreoffice_profile'
                user_profile_dir.mkdir(exist_ok=True)

                try:
                    # Конвертируем в .pptx
                    result = subprocess.run(
                        [
                            soffice_cmd,
                            '--headless',
                            '--convert-to', 'pptx',
                            '--outdir', str(temp_dir_path),
                            '-env:UserInstallation=file:///' + str(user_profile_dir).replace('\\', '/'),
                            str(temp_input)
                        ],
                        capture_output=True,
                        text=True,
                        timeout=120  # 2 минуты таймаут
                    )

                    if result.returncode != 0:
                        error_msg = f"LibreOffice ошибка при конвертации {input_path.name}"
                        if result.stderr:
                            logger.error(f"{error_msg}: {result.stderr}")
                        else:
                            logger.error(error_msg)
                        return False

                except subprocess.TimeoutExpired:
                    logger.error(f"LibreOffice таймаут при конвертации {input_path.name}")
                    return False

                # Ищем созданный .pptx файл
                pptx_file = temp_dir_path / f"{input_path.stem}.pptx"

                if result.stdout:
                    logger.debug(f"LibreOffice stdout: {result.stdout}")

                if not pptx_file.exists():
                    # Иногда LibreOffice создает файл с другим именем
                    pptx_files = list(temp_dir_path.glob('*.pptx'))
                    if pptx_files:
                        pptx_file = pptx_files[0]
                        logger.debug(f"LibreOffice создал файл: {pptx_file.name}")
                    else:
                        logger.error(f"LibreOffice не создал .pptx файл: {pptx_file}")

                        # Детальная диагностика
                        all_files = list(temp_dir_path.iterdir())
                        logger.error(f"Файлы в temp директории: {[f.name for f in all_files]}")

                        if result.stdout:
                            logger.error(f"LibreOffice stdout: {result.stdout}")
                        if result.stderr:
                            logger.error(f"LibreOffice stderr: {result.stderr}")

                        logger.info(f"Для диагностики запустите: python check_libreoffice.py")
                        return False

                # Шаг 2: Конвертируем полученный .pptx в markdown через python-pptx
                logger.debug(f"Конвертация .pptx → .md через python-pptx: {pptx_file.name}")

                success = self._convert_with_pptx(pptx_file, output_path)

                if success:
                    logger.success(f"Успешно сконвертирован .ppt файл: {input_path.name}")
                    return True
                else:
                    logger.error(f"Не удалось сконвертировать промежуточный .pptx файл")
                    return False

            except Exception as e:
                logger.error(f"Ошибка LibreOffice конвертации {input_path.name}: {e}")
                return False
